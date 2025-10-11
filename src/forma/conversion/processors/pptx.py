"""Processor for PPTX files."""

from __future__ import annotations

import logging
import os
import time
import tempfile
from pathlib import Path
from typing import List, Optional


from .base import Processor, ProcessingResult
from ...ocr import parse_image_to_markdown, AdvancedOCRClient
from ...vision import VlmParser, VLMClient
from ...shared.utils.converters import convert_to_pdf
from ...shared.config import get_ocr_config
from ...shared.utils.retry import retry
import fitz  # PyMuPDF


logger = logging.getLogger(__name__)


class PptxProcessor(Processor):
    """Processor for PPTX files with slide-wise heuristics."""

    COMPLEX_THRESHOLD = 25
    
    def __init__(self, vlm_client: Optional[VLMClient] = None, advanced_ocr_client: Optional[AdvancedOCRClient] = None, use_advanced_ocr: bool = False) -> None:
        """Initialize the PPTX processor.
        
        Args:
            vlm_client: Optional VLM client for image description
            advanced_ocr_client: Optional advanced OCR client for text recognition
        """
        self._vlm_client = vlm_client
        self._use_advanced_ocr = use_advanced_ocr
        
        # 仅在开启开关时才尝试创建高级OCR客户端
        self._advanced_ocr_client = None
        if self._use_advanced_ocr:
            self._advanced_ocr_client = advanced_ocr_client
            if self._advanced_ocr_client is None:
                try:
                    config = get_ocr_config()
                    self._advanced_ocr_client = AdvancedOCRClient(
                        api_key=config.api_key,
                        model=config.model,
                        base_url=config.base_url,
                        max_file_size=config.max_file_size
                    )
                    logger.debug(
                        "PptxProcessor: Advanced OCR client initialized with model %s",
                        config.model,
                    )
                except Exception as e:
                    logger.warning(
                        "PptxProcessor: Failed to initialize Advanced OCR client: %s",
                        e,
                    )
                    self._advanced_ocr_client = None

    def _parse_with_retry(self, image_path: Path, slide_idx: int) -> str:
        """使用重试机制调用VLM服务解析幻灯片"""
        
        @retry(
            max_tries=3,
            delay=1.0,
            backoff=2.0,
            exceptions=(Exception,),
            on_retry=lambda e, i: logger.warning(
                "PptxProcessor: VLM retry %s/3 for slide %s due to: %s",
                i,
                slide_idx + 1,
                e,
            ),
        )
        def _parse(path):
            from ...vision import VlmParser
            vlm = VlmParser(self._vlm_client)
            return vlm.parse(path)
        
        return _parse(image_path)
    
    def _recognize_text_with_retry(self, image_path: Path, slide_idx: int) -> str:
        """使用重试机制调用高级OCR服务识别图片文字"""
        
        @retry(
            max_tries=3,
            delay=1.0,
            backoff=2.0,
            exceptions=(Exception,),
            on_retry=lambda e, i: logger.warning(
                "PptxProcessor: GOT-OCR2_0 retry %s/3 for slide %s due to: %s",
                i,
                slide_idx + 1,
                e,
            ),
        )
        def _recognize(path):
            return self._advanced_ocr_client.recognize_text(path)
        
        return _recognize(image_path)
    
    def process(self, input_path: Path) -> ProcessingResult:
        """处理PPTX文件，返回处理结果"""
        process_start_time = time.time()
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE

        # MSO_SHAPE_TYPE.SMART_ART is 24, but might not be in older pptx versions
        MSO_SHAPE_TYPE_SMART_ART = getattr(MSO_SHAPE_TYPE, "SMART_ART", 24)

        path = Path(input_path)
        pres = Presentation(str(path))
        slide_count = len(pres.slides)
        # placeholders for ordered markdown output
        slide_markdowns: List[str] = ["" for _ in range(slide_count)]
        complex_slide_texts: dict[int, str] = {}
        complex_indices: List[int] = []
        image_count = 0

        with tempfile.TemporaryDirectory() as tmpdir:
            tmp = Path(tmpdir)
            # First pass: extract content and decide slide complexity
            for idx, slide in enumerate(pres.slides):
                texts: List[str] = []
                images: List[Path] = []
                is_complex_shape_found = False

                for shape in slide.shapes:
                    # Rule 1: Detect complex shapes that require visual rendering
                    if shape.shape_type in [
                        MSO_SHAPE_TYPE.CHART,
                        MSO_SHAPE_TYPE_SMART_ART,
                        MSO_SHAPE_TYPE.TABLE,  # Tables can be complex
                    ]:
                        is_complex_shape_found = True
                        break

                    if getattr(shape, "has_text_frame", False) and shape.text.strip():
                        texts.append(shape.text.strip())

                    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                        image = shape.image
                        ext = image.ext or "png"
                        img_path = tmp / f"slide{idx}_{len(images)}.{ext}"
                        img_path.write_bytes(image.blob)
                        images.append(img_path)

                slide_text = "\n".join(texts)
                char_count = len(slide_text.replace("\n", "").strip())

                # Rule 2: Slides with very little text are also considered complex
                is_complex = is_complex_shape_found or (char_count < self.COMPLEX_THRESHOLD)

                if is_complex:
                    complex_indices.append(idx)
                    # Store the extracted text for later merging
                    complex_slide_texts[idx] = slide_text
                    # We still process images on complex slides for OCR, just in case
                    # VLM fails or for a more complete picture.

                # All slides (simple and complex) can have images that need OCR
                ocr_texts = [parse_image_to_markdown(str(p)) for p in images]
                if ocr_texts:
                    slide_text = (
                        slide_text + "\n\n" + "\n\n".join(ocr_texts)
                    ).strip()
                
                if not is_complex:
                    # This is a content slide, process with fast path
                    slide_markdowns[idx] = slide_text
                
                image_count += len(images)

            # Deep path for complex slides, optimized for single conversion
            if complex_indices:
                pdf_path = None
                try:
                    # Convert the entire PPTX to PDF once
                    pdf_path = convert_to_pdf(input_path=path, output_dir=tmp)
                    doc = fitz.open(str(pdf_path))
                    vlm = VlmParser(self._vlm_client)

                    for idx in complex_indices:
                        if not (0 <= idx < len(doc)):
                            slide_markdowns[idx] = f"> _[警告] 幻灯片索引 {idx + 1} 超出范围。_"
                            continue
                        
                        page = doc.load_page(idx)
                        pix = page.get_pixmap(dpi=200)
                        img_path = tmp / f"complex_slide_{idx}.png"
                        pix.save(str(img_path))
                        
                        # 深度路径：VLM优先，其次（开启时）高级OCR，最后快速路径文本兜底
                        final_markdown = ""
                        file_size = os.path.getsize(img_path)

                        # 1) VLM 优先
                        if self._vlm_client:
                            try:
                                logger.debug(
                                    "PptxProcessor: Processing slide %s with VLM",
                                    idx + 1,
                                )
                                vlm_markdown = vlm.parse(img_path)
                                if vlm_markdown.strip():
                                    final_markdown = vlm_markdown
                                    logger.debug(
                                        "PptxProcessor: VLM completed for slide %s, text length: %s",
                                        idx + 1,
                                        len(vlm_markdown),
                                    )
                                else:
                                    logger.debug(
                                        "PptxProcessor: VLM returned empty result for slide %s",
                                        idx + 1,
                                    )
                            except Exception as e:
                                logger.error(
                                    "PptxProcessor: VLM failed for slide %s: %s",
                                    idx + 1,
                                    e,
                                )

                        # 2) 若VLM失败或为空，且开启高级OCR，则尝试高级OCR
                        if not final_markdown and self._use_advanced_ocr and self._advanced_ocr_client:
                            try:
                                logger.debug(
                                    "PptxProcessor: Processing slide %s with GOT-OCR2_0 (file size: %s bytes)",
                                    idx + 1,
                                    file_size,
                                )
                                ocr_text = self._advanced_ocr_client.recognize_text(img_path)
                                if ocr_text.strip():
                                    final_markdown = ocr_text
                                    logger.debug(
                                        "PptxProcessor: GOT-OCR2_0 completed for slide %s, text length: %s",
                                        idx + 1,
                                        len(ocr_text),
                                    )
                                else:
                                    logger.debug(
                                        "PptxProcessor: GOT-OCR2_0 returned empty result for slide %s",
                                        idx + 1,
                                    )
                            except ValueError as e:
                                logger.warning(
                                    "PptxProcessor: GOT-OCR2_0 skipped for slide %s: %s",
                                    idx + 1,
                                    e,
                                )
                            except Exception as e:
                                logger.error(
                                    "PptxProcessor: GOT-OCR2_0 failed for slide %s: %s",
                                    idx + 1,
                                    e,
                                )
                        # 3) 快速路径文本兜底
                        if not final_markdown:
                            fast_path_text = complex_slide_texts.get(idx, "")
                            if fast_path_text:
                                final_markdown = fast_path_text
                                logger.debug(
                                    "PptxProcessor: Using fast-path text for slide %s, length: %s",
                                    idx + 1,
                                    len(fast_path_text),
                                )
                            else:
                                final_markdown = "> _[提示] 无法解析此幻灯片内容。_"
                                logger.warning(
                                    "PptxProcessor: No content extracted for slide %s",
                                    idx + 1,
                                )
                        else:
                            # 如果有快速路径文本，并且与最终结果不重复，则合并
                            fast_path_text = complex_slide_texts.get(idx, "")
                            if fast_path_text and fast_path_text not in final_markdown:
                                final_markdown = f"{fast_path_text}\n\n{final_markdown}"

                        slide_markdowns[idx] = final_markdown
                        image_count += 1
                        logger.debug(
                            "PptxProcessor: Completed processing slide %s", idx + 1
                        )

                    doc.close()

                except (RuntimeError, FileNotFoundError) as e:
                    # If conversion or parsing fails, provide a general message
                    # and leave the complex slide markdown empty or with a notice.
                    error_msg = (
                        f"> _[提示] 深度解析失败: {e}_\n"
                        f"> _请确保已正确安装 LibreOffice。_"
                    )
                    logger.error("PptxProcessor: Deep parsing failed: %s", e)
                    for idx in complex_indices:
                        if not slide_markdowns[idx]: # Avoid overwriting fallback text
                            slide_markdowns[idx] = error_msg

        markdown = "\n\n---\n\n".join(m for m in slide_markdowns if m)
        text_len = len(markdown.strip())
        total_elapsed = time.time() - process_start_time
        logger.debug(
            "PptxProcessor: Processing completed in %.2fs", total_elapsed
        )
        
        return ProcessingResult(
            markdown_content=markdown,
            text_char_count=text_len,
            image_count=image_count,
            low_confidence=text_len == 0,
        )


__all__ = ["PptxProcessor"]
