"""Processor for PPTX files."""

from __future__ import annotations

import tempfile
from pathlib import Path
from typing import List

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from .base import Processor, ProcessingResult
from ..ocr import parse_image_to_markdown
from ..vlm import VlmParser
from ...utils.converters import convert_ppt_slide_to_image


class PptxProcessor(Processor):
    """Processor for PPTX files with slide-wise heuristics."""

    COMPLEX_THRESHOLD = 25

    def process(self, input_path: Path) -> ProcessingResult:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE

        # MSO_SHAPE_TYPE.SMART_ART is 24, but might not be in older pptx versions
        MSO_SHAPE_TYPE_SMART_ART = getattr(MSO_SHAPE_TYPE, "SMART_ART", 24)

        path = Path(input_path)
        pres = Presentation(str(path))
        slide_count = len(pres.slides)
        # placeholders for ordered markdown output
        slide_markdowns: List[str] = ["" for _ in range(slide_count)]
        complex_indices: List[int] = []
        image_count = 0

        with tempfile.TemporaryDirectory() as tmpdir:
            tmp = Path(tmpdir)
            # First pass: extract content and decide slide complexity
            for idx, slide in enumerate(pres.slides):
                texts: List[str] = []
                images: List[Path] = []
                is_complex_shape_found = False

                for shape in slide.shapes:
                    # Rule 1: Detect complex shapes that require visual rendering
                    if shape.shape_type in [
                        MSO_SHAPE_TYPE.CHART,
                        MSO_SHAPE_TYPE_SMART_ART,
                        MSO_SHAPE_TYPE.TABLE,  # Tables can be complex
                    ]:
                        is_complex_shape_found = True
                        break

                    if getattr(shape, "has_text_frame", False) and shape.text.strip():
                        texts.append(shape.text.strip())

                    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                        image = shape.image
                        ext = image.ext or "png"
                        img_path = tmp / f"slide{idx}_{len(images)}.{ext}"
                        img_path.write_bytes(image.blob)
                        images.append(img_path)

                # If a complex shape was found, the whole slide is complex
                if is_complex_shape_found:
                    complex_indices.append(idx)
                    continue

                slide_text = "\n".join(texts)
                char_count = len(slide_text.replace("\n", "").strip())

                # Rule 2: Slides with very little text are complex
                if char_count < self.COMPLEX_THRESHOLD:
                    complex_indices.append(idx)
                else:
                    # This is a content slide, process with fast path
                    ocr_texts = [parse_image_to_markdown(
                        str(p)) for p in images]
                    if ocr_texts:
                        slide_text = (
                            slide_text + "\n\n" + "\n\n".join(ocr_texts)
                        ).strip()
                    slide_markdowns[idx] = slide_text
                    image_count += len(images)

            # Deep path for complex slides
            if complex_indices:
                vlm = VlmParser()
                for idx in complex_indices:
                    try:
                        img_path = convert_ppt_slide_to_image(
                            ppt_path=path, slide_index=idx, output_dir=tmp
                        )
                        slide_markdowns[idx] = vlm.parse(img_path)
                        image_count += 1
                    except (RuntimeError, ValueError, FileNotFoundError) as e:
                        slide_markdowns[idx] = (
                            f"> _[提示] 第 {idx + 1} 页是一张复杂幻灯片，深度解析失败：{e}_\n"
                            f"> _系统已回退并保留可提取的文本/图片。若因未安装 LibreOffice 导致，请安装并确保其在系统路径或 macOS 默认位置。_"
                        )

        markdown = "\n\n---\n\n".join(m for m in slide_markdowns if m)
        text_len = len(markdown.strip())
        return ProcessingResult(
            markdown_content=markdown,
            text_char_count=text_len,
            image_count=image_count,
            low_confidence=text_len == 0,
        )


__all__ = ["PptxProcessor"]
