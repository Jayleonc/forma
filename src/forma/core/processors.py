"""Document processors for the fast pipeline."""

from __future__ import annotations

from abc import ABC, abstractmethod
from dataclasses import dataclass
from pathlib import Path
from concurrent.futures import ThreadPoolExecutor, as_completed
from typing import List
import tempfile

from .ocr import parse_image_to_markdown
from .vlm import VlmParser
from ..utils.converters import convert_ppt_slide_to_image


@dataclass
class ProcessingResult:
    """Result returned by processors."""

    markdown_content: str
    text_char_count: int
    image_count: int
    low_confidence: bool = False


class Processor(ABC):
    """Abstract processor for a single document type."""

    @abstractmethod
    def process(self, input_path: Path) -> ProcessingResult:  # pragma: no cover - interface
        """Process the given file and return a :class:`ProcessingResult`."""
        raise NotImplementedError


class PdfProcessor(Processor):
    """Processor for PDF files using PyMuPDF and OCR."""

    def process(self, input_path: Path) -> ProcessingResult:
        import pymupdf4llm
        import fitz

        path = Path(input_path)
        base_md = pymupdf4llm.to_markdown(str(path))
        text_len = len(base_md.strip())

        doc = fitz.open(str(path))
        image_paths: List[Path] = []
        with tempfile.TemporaryDirectory() as tmpdir:
            tmp = Path(tmpdir)
            for page_index, page in enumerate(doc):
                for img_index, img in enumerate(page.get_images(full=True)):
                    xref = img[0]
                    base_image = doc.extract_image(xref)
                    ext = base_image.get("ext", "png")
                    img_bytes = base_image["image"]
                    img_path = tmp / f"p{page_index}_{img_index}.{ext}"
                    img_path.write_bytes(img_bytes)
                    image_paths.append(img_path)
            doc.close()

            ocr_texts: List[str] = []
            if image_paths:
                from . import parser as _parser

                with ThreadPoolExecutor() as executor:
                    futures = [
                        executor.submit(_parser.ocr_image_file, str(p))
                        for p in image_paths
                    ]
                    for future in as_completed(futures):
                        ocr_texts.append(future.result())

        markdown = base_md
        if ocr_texts:
            appendix = (
                "\n\n---\n\n## 附录：图片内容解析\n\n" + "\n\n---\n\n".join(ocr_texts)
            )
            markdown += appendix

        low_conf = text_len < 50
        return ProcessingResult(
            markdown_content=markdown,
            text_char_count=text_len,
            image_count=len(image_paths),
            low_confidence=low_conf,
        )


class ImageProcessor(Processor):
    """Processor for image files using OCR."""

    def process(self, input_path: Path) -> ProcessingResult:
        md = parse_image_to_markdown(str(input_path))
        text_len = len(md.strip())
        return ProcessingResult(
            markdown_content=md,
            text_char_count=text_len,
            image_count=1,
            low_confidence=text_len == 0,
        )


class DocxProcessor(Processor):
    """Processor for DOCX files using a hybrid approach."""

    def process(self, input_path: Path) -> ProcessingResult:
        path = str(input_path)
        md = None

        # Plan B: High-fidelity conversion with Mammoth
        # 把 Docs 转成 HTML，再转成 Markdown（保留表格结构）
        try:
            import mammoth 
            import markdownify 

            with open(path, "rb") as f:
                html = mammoth.convert_to_html(f).value

            # 转换 HTML 到 Markdown，保留表格结构
            md = markdownify.markdownify(html, heading_style="ATX").strip()
        except Exception:
            # 如果转换失败，使用 Plan A
            md = None

        # Plan A: Fallback to pure python-docx for robustness
        # 如果 Plan B 失败，使用 Plan A
        if not md:
            from forma.utils.docx import docx_to_markdown_gfm

            md = docx_to_markdown_gfm(path)

        text_len = len(md.strip())

        # Count images using python-docx (as a basic heuristic)
        # 计算图片数量
        from docx import Document 

        doc = Document(path)
        image_count = 0
        for rel in doc.part._rels.values():
            if "image" in rel.target_ref:
                image_count += 1

        return ProcessingResult(
            markdown_content=md,
            text_char_count=text_len,
            image_count=image_count,
            low_confidence=text_len == 0,
        )


class PptxProcessor(Processor):
    """Processor for PPTX files with slide-wise heuristics."""

    COMPLEX_THRESHOLD = 25

    def process(self, input_path: Path) -> ProcessingResult:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE

        path = Path(input_path)
        pres = Presentation(str(path))
        slide_count = len(pres.slides)

        # placeholders for ordered markdown output
        slide_markdowns: List[str] = ["" for _ in range(slide_count)]
        complex_indices: List[int] = []
        image_count = 0

        with tempfile.TemporaryDirectory() as tmpdir:
            tmp = Path(tmpdir)
            # First pass: extract text/images and decide complexity
            for idx, slide in enumerate(pres.slides):
                texts: List[str] = []
                images: List[Path] = []
                for shape in slide.shapes:
                    if getattr(shape, "has_text_frame", False):
                        text = shape.text.strip()
                        if text:
                            texts.append(text)
                    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                        image = shape.image
                        ext = image.ext or "png"
                        img_path = tmp / f"slide{idx}_{len(images)}.{ext}"
                        img_path.write_bytes(image.blob)
                        images.append(img_path)
                slide_text = "\n".join(texts)
                char_count = len(slide_text.replace("\n", "").strip())
                if char_count < self.COMPLEX_THRESHOLD:
                    complex_indices.append(idx)
                else:
                    ocr_texts = [parse_image_to_markdown(str(p)) for p in images]
                    if ocr_texts:
                        slide_text = (slide_text + "\n\n" + "\n\n".join(ocr_texts)).strip()
                    slide_markdowns[idx] = slide_text
                    image_count += len(images)

            # Deep path for complex slides via LibreOffice + VLM
            if complex_indices:
                vlm = VlmParser()
                for idx in complex_indices:
                    try:
                        img_path = convert_ppt_slide_to_image(
                            ppt_path=path, slide_index=idx, output_dir=tmp
                        )
                        slide_markdowns[idx] = vlm.parse(img_path)
                        image_count += 1
                    except (RuntimeError, ValueError) as e:
                        # If conversion fails, add an error message to the markdown.
                        slide_markdowns[idx] = f"_Error processing complex slide {idx + 1}: {e}_"

        markdown = "\n\n---\n\n".join(m for m in slide_markdowns if m)
        text_len = len(markdown.strip())
        return ProcessingResult(
            markdown_content=markdown,
            text_char_count=text_len,
            image_count=image_count,
            low_confidence=text_len == 0,
        )


__all__ = [
    "ProcessingResult",
    "Processor",
    "PdfProcessor",
    "ImageProcessor",
    "DocxProcessor",
    "PptxProcessor",
]
