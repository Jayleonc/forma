from __future__ import annotations

from pathlib import Path
import types

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches
from PIL import Image, ImageDraw

from forma.core.processors import PptxProcessor


def _create_demo_pptx(path: Path) -> None:
    prs = Presentation()
    title_content_layout = prs.slide_layouts[1]

    # Slide 0: Fast path with plenty of text and one image
    slide_fast = prs.slides.add_slide(prs.slide_layouts[5])  # blank
    tb = slide_fast.shapes.add_textbox(
        Inches(1), Inches(1), Inches(8), Inches(2))
    tb.text_frame.text = (
        "This slide has enough textual content to avoid the deep path."
    )

    img = Image.new("RGB", (50, 50), color="white")
    draw = ImageDraw.Draw(img)
    draw.text((5, 5), "hi", fill="black")
    img_path = path.with_suffix(".png")
    img.save(img_path)
    slide_fast.shapes.add_picture(str(img_path), Inches(1), Inches(2))
    img_path.unlink()

    # Slide 1: Complex path with very little text (original rule)
    slide_complex_text = prs.slides.add_slide(prs.slide_layouts[5])
    tb2 = slide_complex_text.shapes.add_textbox(
        Inches(1), Inches(1), Inches(4), Inches(1)
    )
    tb2.text_frame.text = "Hi"

    # Slide 2: Complex path due to containing a table
    slide_table = prs.slides.add_slide(title_content_layout)
    slide_table.shapes.title.text = "Table Slide"
    rows, cols = 2, 2
    table = slide_table.shapes.add_table(
        rows, cols, Inches(2), Inches(2), Inches(4), Inches(1.5)
    ).table
    table.cell(0, 0).text = "A"
    table.cell(0, 1).text = "B"
    table.cell(1, 0).text = "C"
    table.cell(1, 1).text = "D"

    # Slide 3: Complex path due to containing a chart
    slide_chart = prs.slides.add_slide(title_content_layout)
    slide_chart.shapes.title.text = "Chart Slide"
    chart_data = CategoryChartData()
    chart_data.categories = ["East", "West", "Midwest"]
    chart_data.add_series("Series 1", (19.2, 21.4, 16.7))
    slide_chart.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(2), Inches(
            2), Inches(6), Inches(4.5), chart_data
    )

    prs.save(path)


def test_pptx_processor(monkeypatch, tmp_path):
    pptx_path = tmp_path / "demo.pptx"
    _create_demo_pptx(pptx_path)

    # Stub image OCR
    monkeypatch.setattr(
        "forma.core.processors.pptx.parse_image_to_markdown", lambda p: "image md"
    )

    # Stub VLM parser
    calls = {"vlm": 0}

    class DummyVlm:
        def parse(self, path, prompt_name: str = "default_image_description"):
            calls["vlm"] += 1
            return "vlm slide"

    monkeypatch.setattr("forma.core.processors.pptx.VlmParser", DummyVlm)

    # Stub the new converter utility
    def fake_converter(ppt_path: Path, slide_index: int, output_dir: Path) -> Path:
        """A mock for the converter that creates a dummy image."""
        img_path = output_dir / f"complex_slide_{slide_index}.png"
        Image.new("RGB", (10, 10), color="blue").save(img_path)
        return img_path

    monkeypatch.setattr(
        "forma.core.processors.pptx.convert_ppt_slide_to_image", fake_converter
    )

    processor = PptxProcessor()
    result = processor.process(pptx_path)

    assert "image md" in result.markdown_content
    assert "vlm slide" in result.markdown_content
    # We expect 3 slides to be complex: little text, table, and chart
    assert calls["vlm"] == 3
